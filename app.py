from dotenv import load_dotenv
import streamlit as st
from PyPDF2 import PdfReader
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
from langchain.callbacks import get_openai_callback
from docx import Document
from pptx import Presentation
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from transformers import pipeline, BartTokenizer, BartForConditionalGeneration
from scipy.cluster import hierarchy
import openpyxl
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_community.chat_models import ChatOpenAI

def extract_column_name(prompt):
    """Extract column name from a given prompt."""
    keywords = ["of", "on", "from"]
    words = prompt.split()
    for i, word in enumerate(words):
        if word in keywords and i < len(words) - 1:
            return words[i + 1]
    return None

def extract_person_name(query):
    """Extract person name from a given query."""
    words = query.split()
    for i, word in enumerate(words):
        if word.lower() == "name" and i < len(words) - 1:
            return words[i + 1]
    return None

def summarize_text(summarizer, text, max_chunk_length=512):
    """Summarize the provided text using the given summarizer model."""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=max_chunk_length,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    summaries = []
    for chunk in chunks:
        chunk_length = len(chunk.split())
        max_length = min(250, max(chunk_length // 2, 50))  # Adjust max_length based on chunk size
        summary = summarizer(chunk, max_length=max_length, min_length=30, do_sample=False)
        summaries.append(summary[0]['summary_text'])
    return "\n".join(summaries)

def generate_quiz(quiz_generator, text, max_chunk_length=512):
    """Generate quiz questions from the provided text using a text2text-generation model."""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=max_chunk_length,
        chunk_overlap=200,
        length_function=len
    )
    chunks = text_splitter.split_text(text)
    quizzes = []
    for chunk in chunks:
        input_text = f"generate question: {chunk}"
        quiz = quiz_generator(input_text, max_length=50, num_return_sequences=1)
        quizzes.append(quiz)
    return quizzes

def translate_text(translator, text, max_length=400):
    """Translate the provided text to the target language using the given translator model."""
    chunks = [text[i:i + max_length] for i in range(0, len(text), max_length)]
    translations = []
    for chunk in chunks:
        translated = translator(chunk, max_length=max_length)
        translations.append(translated[0]['translation_text'])
    return " ".join(translations)

def main():
    load_dotenv()
    st.set_page_config(
        page_title="DocuWiz AI",
        page_icon="📚",
        initial_sidebar_state="expanded",
    )
   
    st.markdown(
        """
        <style>
        .stApp header {
            background-color: #4CAF50;
        }
        .stApp header a, .stApp header h1, .stApp header h2 {
            color: white;
        }
        .main-header {
            text-align: center;
            font-size: 2.5rem;
            color: #4CAF50;
            margin-bottom: 2rem;
        }
        .sub-header {
            text-align: center;
            font-size: 1.2rem;
            color: #555;
            margin-bottom: 2rem;
        }
        .contact-header {
            text-align: center;
            font-size: 2rem;
            color: #4CAF50;
            margin-bottom: 1rem;
        }
        .contact-info {
            font-size: 1.1rem;
            margin-bottom: 1rem;
        }
        </style>
        """,
        unsafe_allow_html=True,
    )

    # Add sidebar for navigation
    st.sidebar.title("Navigation")
    page = st.sidebar.radio("Go to", ["Home", "Contact Us"])

    if page == "Home":
        st.markdown("<h1 class='main-header'>DocuWiz AI 📚</h1>", unsafe_allow_html=True)
        st.markdown("<h3 class='sub-header'>Empowering Document Intelligence with AI</h3>", unsafe_allow_html=True)
        st.markdown("**This application is powered by Jillani SoftTech 😎**")

        uploaded_file = st.file_uploader("Upload your document", type=["pdf", "docx", "txt", "pptx", "csv", "xlsx"], key="file_uploader", help="Supported file types: PDF, DOCX, TXT, PPTX, CSV, XLSX", accept_multiple_files=False)

        if uploaded_file is not None:
            text = extract_text_from_file(uploaded_file)
            
            if text:
                model = BartForConditionalGeneration.from_pretrained("facebook/bart-large", forced_bos_token_id=0)
                #tok = BartTokenizer.from_pretrained("facebook/bart-large")
                tokenizer = BartTokenizer.from_pretrained("facebook/bart-large")
                max_token_length = 1024  # Increased to allow longer inputs
                tokens = tokenizer.tokenize(text)
                if len(tokens) > max_token_length:
                    text = tokenizer.convert_tokens_to_string(tokens[:max_token_length])

                text_splitter = RecursiveCharacterTextSplitter(
                    chunk_size=max_token_length,
                    chunk_overlap=200,
                    length_function=len
                )
                chunks = text_splitter.split_text(text)

                embeddings = OpenAIEmbeddings(model="text-embedding-3-small")
                knowledge_base = FAISS.from_texts(chunks, embedding=embeddings)

                user_question = st.text_input("Ask a question about your document:")
                if user_question:
                    docs = knowledge_base.similarity_search(user_question)

                    llm = ChatOpenAI(model_name="gpt-4o", temperature=0.2,max_tokens=1000)
                    chain = load_qa_chain(llm, chain_type="stuff")
                    with get_openai_callback() as cb:
                        response = chain.run(input_documents=docs, question=user_question)

                    st.write(response)

                # Initialize NLP pipelines
                summarizer = pipeline("summarization", model="facebook/bart-large-cnn") 
                sentiment_analyzer = pipeline("sentiment-analysis", model="distilbert-base-uncased-finetuned-sst-2-english")
                topic_modeler = pipeline("zero-shot-classification", model="facebook/bart-large-mnli")
                quiz_generator = pipeline("text2text-generation", model="valhalla/t5-small-e2e-qg")
                translator_models = {
                    "ur": "Helsinki-NLP/opus-mt-en-ur",
                    "fr": "Helsinki-NLP/opus-mt-en-fr",
                    "es": "Helsinki-NLP/opus-mt-en-es",
                    "de": "Helsinki-NLP/opus-mt-en-de",
                    "it": "Helsinki-NLP/opus-mt-en-it",
                    "zh": "Helsinki-NLP/opus-mt-en-zh"
                }

                if st.button("Summarize Document"):
                    try:
                        if len(text.strip()) == 0:
                            st.write("The document is too short or empty.") 
                        else:
                            summary = summarize_text(summarizer, text)
                            st.write(summary)
                    except Exception as e:
                        st.error(f"Error during summarization: {str(e)}")

                if st.button("Analyze Sentiment"):
                    if len(text) == 0:
                        st.write("The document is empty.")
                    else:
                        sentiment = sentiment_analyzer(text)
                        st.write(f"Sentiment: {sentiment[0]['label']} (Score: {sentiment[0]['score']})")

                if st.button("Extract Topics"):
                    if len(text) == 0:
                        st.write("The document is empty.")
                    else:
                        labels = ["Politics", "Economics", "Sports", "Technology", "Health", "Entertainment", "Education"]
                        topics = topic_modeler(text, candidate_labels=labels)
                        st.write(f"Main Topic: {topics['labels'][0]} (Score: {topics['scores'][0]})")

                # Adding translation feature
                target_lang = st.selectbox("Select target language", ["ur","fr", "es", "de", "it", "zh"])
                if target_lang and st.button("Translate Document"):
                    translator_model = translator_models[target_lang]
                    translator = pipeline("translation", model=translator_model)
                    try:
                        translation = translate_text(translator, text)
                        st.write(translation)
                    except Exception as e:
                        st.error(f"Error during translation: {str(e)}")

                # Adding quiz generation feature
                if st.button("Generate Quiz"):
                    try:
                        quizzes = generate_quiz(quiz_generator, text)
                        for i, quiz in enumerate(quizzes):
                            st.write(f"Question {i+1}: {quiz[0]['generated_text']}")
                    except Exception as e:
                        st.error(f"Error during quiz generation: {str(e)}")

            else:
                st.warning("Unsupported file type. Please upload a PDF, DOCX, TXT, PPTX, XLSX, or CSV file.")
            
        # Additional handling for CSV and Excel files with data analysis features
        handle_csv_excel_files(uploaded_file)

    elif page == "Contact Us":
        st.markdown("<h1 class='contact-header'>Contact Us</h1>", unsafe_allow_html=True)
        st.markdown(
            """
            <div class='contact-info'>
                <p><strong>Email:</strong> <a href='mailto:m.g.jillani123@gmail.com'>m.g.jillani123@gmail.com</a></p>
                <p><strong>Contact number:</strong> +92-321-1174167</p>
                <p><strong>LinkedIn:</strong> <a href='https://linkedin.com/in/jillanisofttech' target='_blank'>linkedin.com/in/jillanisofttech</a></p>
                <p><strong>Website:</strong> <a href='https://mgjillanimughal.github.io' target='_blank'>mgjillanimughal.github.io</a></p>
                <p><strong>Kaggle:</strong> <a href='https://kaggle.com/jillanisofttech' target='_blank'>kaggle.com/jillanisofttech</a></p>
                <p><strong>GitHub:</strong> <a href='https://github.com/MGJillaniMughal' target='_blank'>github.com/MGJillaniMughal</a></p>
                <p><strong>Medium:</strong> <a href='https://jillanisofttech.medium.com' target='_blank'>jillanisofttech.medium.com</a></p>
                <p><strong>Calendly:</strong> <a href='https://calendly.com/jillanisofttech/30mins' target='_blank'>calendly.com/jillanisofttech/30mins</a></p>
            </div>
            """,
            unsafe_allow_html=True
        )

def extract_text_from_file(uploaded_file):
    """Extract text from different file types."""
    try:
        if uploaded_file.type == 'application/pdf':
            pdf_reader = PdfReader(uploaded_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
            return text
        elif uploaded_file.type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            doc = Document(uploaded_file)
            return "\n".join([para.text for para in doc.paragraphs])
        elif uploaded_file.type == 'text/plain':
            return uploaded_file.read().decode('utf-8')
        elif uploaded_file.type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            ppt = Presentation(uploaded_file)
            text = ""
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text'):  
                        text += shape.text + "\n"
            return text
    except Exception as e:
        st.error(f"Error extracting text: {str(e)}")
    return None

def handle_csv_excel_files(uploaded_file):
    """Handle CSV and Excel file uploads with data analysis features."""
    if uploaded_file is not None:
        try:
            if uploaded_file.type == 'text/csv':
                df = pd.read_csv(uploaded_file)
                st.write("Uploaded CSV Data:")
                st.write(df)

                prompt = st.text_input("Enter a data science analysis prompt:")

                if prompt:
                    st.write("Analysis Result:")
                    prompt_lower = prompt.lower()
                    column = extract_column_name(prompt)
                    if column is not None:
                        perform_data_analysis(df, column, prompt_lower)
                    else:
                        st.write("Column not specified in the prompt.")
                return
            elif uploaded_file.type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                df = pd.read_excel(uploaded_file, engine="openpyxl")
                st.write("Uploaded Excel Data:")
                st.write(df)

                user_prompt = st.text_input("Enter a prompt for Excel analysis:")

                if user_prompt:
                    st.write("Excel Analysis Result:")
                    user_prompt_lower = user_prompt.lower()
                    column = extract_column_name(user_prompt)
                    if column is not None:
                        perform_data_analysis(df, column, user_prompt_lower)
                    else:
                        st.write("Column not specified in the prompt.")
                return
        except Exception as e:
            st.error(f"Error processing file: {str(e)}")

def perform_data_analysis(df, column, prompt_lower):
    """Perform various data analysis tasks based on the user prompt."""
    try:
        if "mean" in prompt_lower or "average" in prompt_lower:
            st.write(f"Mean of '{column}': {df[column].mean()}")

        elif "median" in prompt_lower:
            st.write(f"Median of '{column}': {df[column].median()}")

        elif "mode" in prompt_lower:
            mode_values = df[column].mode()
            if not mode_values.empty:
                st.write(f"Mode of '{column}': {', '.join(map(str, mode_values))}")
            else:
                st.write(f"No mode found in '{column}'.")

        elif "histogram" in prompt_lower:
            plt.hist(df[column], bins=20)
            st.pyplot(plt)

        elif "scatterplot" in prompt_lower:
            x_column = st.selectbox("Select the X-axis column:", df.columns)
            plt.scatter(df[x_column], df[column])
            st.pyplot(plt)

        elif "count" in prompt_lower:
            st.write(f"Count of '{column}': {df[column].count()}")

        elif "sum" in prompt_lower:
            st.write(f"Sum of '{column}': {df[column].sum()}")

        elif "null" in prompt_lower:
            st.write(f"Null value count in '{column}': {df[column].isnull().sum()}")

        elif "min" in prompt_lower:
            st.write(f"Min value in '{column}': {df[column].min()}")

        elif "max" in prompt_lower:
            st.write(f"Max value in '{column}': {df[column].max()}")

        elif "line plot" in prompt_lower:
            x_column = st.selectbox("Select the X-axis column:", df.columns)
            y_column = st.selectbox("Select the Y-axis column:", df.columns)
            plt.plot(df[x_column], df[y_column])
            st.pyplot(plt)

        elif "scatter chart" in prompt_lower:
            x_column = st.selectbox("Select the X-axis column:", df.columns)
            y_column = st.selectbox("Select the Y-axis column:", df.columns)
            plt.scatter(df[x_column], df[y_column])
            st.pyplot(plt)

        elif "correlation chart" in prompt_lower:
            corr_matrix = df.corr()
            sns.heatmap(corr_matrix, annot=True)
            st.pyplot(plt)

        elif "heatmap" in prompt_lower:
            sns.heatmap(df.corr(), annot=True, cmap="coolwarm")
            st.pyplot(plt)

        elif "bubble chart" in prompt_lower:
            x_column = st.selectbox("Select the X-axis column:", df.columns)
            y_column = st.selectbox("Select the Y-axis column:", df.columns)
            size_column = st.selectbox("Select the size column:", df.columns)
            plt.scatter(df[x_column], df[y_column], s=df[size_column])
            st.pyplot(plt)

        elif "radar chart" in prompt_lower:
            st.write("Radar chart not implemented yet.")

        elif "ridge plot" in prompt_lower:
            sns.violinplot(df[column])
            st.pyplot(plt)

        elif "dendrogram" in prompt_lower:
            corr_matrix = df.corr()
            linkage_matrix = hierarchy.linkage(corr_matrix, method='ward')
            dendrogram = hierarchy.dendrogram(linkage_matrix, labels=corr_matrix.index)
            st.pyplot(plt)

        else:
            st.write("Unsupported analysis prompt.")
    except Exception as e:
        st.error(f"Error during data analysis: {str(e)}")

if __name__ == '__main__':
    main()
